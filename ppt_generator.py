"""
ppt_generator.py
-----------------
Core logic that reads the LinkedIn-activity Excel file and fills the
"Sample.pptx" template with the data, producing a ready-to-use .pptx.

This file has NO Streamlit code in it on purpose - it can be imported
and tested on its own (see test_generate.py) as well as used by app.py.
"""

import copy
import re
from io import BytesIO

import openpyxl
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# 1. Settings you may want to tweak
# ---------------------------------------------------------------------------

# Column headers expected in the Excel file (row 1).
EXCEL_COL_POST_TYPE = "PostType"
EXCEL_COL_TIME_AGO = "TimeAgo"
EXCEL_COL_CONTENT = "Content"
EXCEL_COL_POST_LINK = "PostLink"
EXCEL_COL_SOURCE = "Source"      # the cell that holds the hyperlink
EXCEL_COL_HEADLINE = "Headline"

# Color used for the "Source" hyperlink text (typical PowerPoint link blue)
HYPERLINK_COLOR = RGBColor(0x05, 0x66, 0xC2)

# Names of the branding shapes in Sample.pptx (Picture / LinkedIn ID / Name /
# Country). These are matched by shape name first; if not found, we fall
# back to matching the template's current placeholder text so this still
# works if someone renames the shapes.
BRAND_PICTURE_NAME = "Picture 4"
BRAND_LINKEDIN_ID_PLACEHOLDER_TEXT = "xxxx/"     # value cell inside "Group 10"
BRAND_NAME_PLACEHOLDER_TEXT = "xxxxxx"           # executive's name, below photo
BRAND_COUNTRY_PLACEHOLDER_TEXT = "xxxx"          # executive's country, below photo


# ---------------------------------------------------------------------------
# 2. Read the Excel file
# ---------------------------------------------------------------------------
def read_excel_rows(file_like):
    """
    Reads the uploaded Excel file and returns a list of dicts, one per row:
        {
          "post_type": ...,
          "timeline":  ...,
          "headline":  ...,
          "content":   ...,
          "source_url": ...,
        }
    The hyperlink is read from the "Source" column cell itself
    (Insert > Link in Excel). If that cell has no hyperlink, we fall back
    to the "PostLink" column's plain text value.
    """
    wb = openpyxl.load_workbook(file_like, data_only=True)
    ws = wb.active

    # Map header name -> column index (1-based), from row 1
    headers = {}
    for cell in ws[1]:
        if cell.value:
            headers[str(cell.value).strip()] = cell.column

    required = [EXCEL_COL_POST_TYPE, EXCEL_COL_TIME_AGO, EXCEL_COL_CONTENT,
                EXCEL_COL_SOURCE, EXCEL_COL_HEADLINE]
    missing = [c for c in required if c not in headers]
    if missing:
        raise ValueError(
            f"The Excel file is missing these expected column headers: {missing}. "
            f"Found headers: {list(headers.keys())}"
        )

    rows = []
    for r in range(2, ws.max_row + 1):
        post_type = ws.cell(row=r, column=headers[EXCEL_COL_POST_TYPE]).value
        timeline = ws.cell(row=r, column=headers[EXCEL_COL_TIME_AGO]).value
        content = ws.cell(row=r, column=headers[EXCEL_COL_CONTENT]).value
        headline = ws.cell(row=r, column=headers[EXCEL_COL_HEADLINE]).value

        # Skip fully blank rows
        if not any([post_type, timeline, content, headline]):
            continue

        source_cell = ws.cell(row=r, column=headers[EXCEL_COL_SOURCE])
        source_url = source_cell.hyperlink.target if source_cell.hyperlink else None
        if not source_url and EXCEL_COL_POST_LINK in headers:
            source_url = ws.cell(row=r, column=headers[EXCEL_COL_POST_LINK]).value

        rows.append({
            "post_type": post_type or "",
            "timeline": timeline or "",
            "headline": headline or "",
            "content": content or "",
            "source_url": source_url or "",
        })

    return rows


# ---------------------------------------------------------------------------
# 3. Slide-duplication helper (python-pptx has no built-in "duplicate slide")
# ---------------------------------------------------------------------------
def duplicate_slide(prs, slide_index):
    """
    Duplicates the slide at slide_index (0-based) and appends the copy at
    the end of the presentation. Returns the new slide.
    Handles pictures / relationships so images are not lost.
    """
    source = prs.slides[slide_index]
    dest = prs.slides.add_slide(source.slide_layout)

    # The new slide auto-inherits placeholder shapes from the layout -
    # remove them, we will copy the real shapes from the source slide instead.
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    # Copy every relationship (images etc.) from source slide to destination,
    # building a map of old-rId -> new-rId so we can fix references below.
    rid_map = {}
    for rel_id, rel in source.part.rels.items():
        if "slideLayout" in rel.reltype or "notesSlide" in rel.reltype:
            continue
        if rel.is_external:
            new_rel_id = dest.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rel_id = dest.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rel_id] = new_rel_id

    # Deep-copy every shape and fix up any r:id / r:embed references
    for shp in source.shapes:
        new_el = copy.deepcopy(shp._element)
        for el in new_el.iter():
            for attr in el.attrib:
                if el.attrib[attr] in rid_map:
                    el.attrib[attr] = rid_map[el.attrib[attr]]
        dest.shapes._spTree.append(new_el)

    return dest


# ---------------------------------------------------------------------------
# 4. Table helpers
# ---------------------------------------------------------------------------
def find_table_shape(slide):
    for shp in slide.shapes:
        if shp.has_table:
            return shp
    raise ValueError("No table found on the template slide.")


def iter_all_shapes(shapes):
    """Yield every shape, descending into groups."""
    for shp in shapes:
        yield shp
        if shp.shape_type == 6:  # GROUP
            yield from iter_all_shapes(shp.shapes)


def find_shape_by_name(slide, name):
    for shp in iter_all_shapes(slide.shapes):
        if shp.name == name:
            return shp
    return None


def find_shape_by_exact_text(slide, text):
    for shp in iter_all_shapes(slide.shapes):
        if shp.has_text_frame and shp.text_frame.text.strip() == text:
            return shp
    return None


def find_picture_shape(slide, name):
    for shp in iter_all_shapes(slide.shapes):
        if shp.shape_type == 13 and shp.name == name:  # PICTURE
            return shp
    # fallback: first picture shape found anywhere on the slide
    for shp in iter_all_shapes(slide.shapes):
        if shp.shape_type == 13:
            return shp
    return None


def replace_picture(slide, picture_shape, image_bytes):
    """Swap the image inside an existing picture shape, keeping its
    position, size and crop exactly as in the template."""
    image_part, rId = slide.part.get_or_add_image_part(BytesIO(image_bytes))
    blip = picture_shape._element.blipFill.blip
    blip.rEmbed = rId


def set_branding(slide, linkedin_id=None, exec_name=None, exec_country=None,
                  picture_bytes=None):
    """Fills in the executive-specific header info: LinkedIn ID, name,
    country and profile photo. Any argument left as None/empty is skipped
    (the template's existing value stays)."""
    if linkedin_id:
        shp = (find_shape_by_exact_text(slide, BRAND_LINKEDIN_ID_PLACEHOLDER_TEXT)
               or find_shape_by_exact_text(slide, "xxxxxxxxxxxxxxxxxxxxxxxx/"))
        if shp is not None:
            set_cell_text_frame(shp.text_frame, linkedin_id)

    if exec_name:
        shp = find_shape_by_exact_text(slide, BRAND_NAME_PLACEHOLDER_TEXT)
        if shp is not None:
            set_cell_text_frame(shp.text_frame, exec_name)

    if exec_country:
        shp = find_shape_by_exact_text(slide, BRAND_COUNTRY_PLACEHOLDER_TEXT)
        if shp is not None:
            set_cell_text_frame(shp.text_frame, exec_country)

    if picture_bytes:
        pic_shape = find_picture_shape(slide, BRAND_PICTURE_NAME)
        if pic_shape is not None:
            replace_picture(slide, pic_shape, picture_bytes)


def set_cell_text_frame(text_frame, text):
    """Like set_cell_text but works on any text_frame (not just table cells)."""
    p = text_frame.paragraphs[0]
    if p.runs:
        run = p.runs[0]
        run.text = str(text)
        for extra in p.runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        run = p.add_run()
        run.text = str(text)
    for extra_p in text_frame.paragraphs[1:]:
        extra_p._p.getparent().remove(extra_p._p)


def find_title_shape(slide):
    for shp in slide.shapes:
        if shp.has_text_frame and shp.name.lower().startswith("title"):
            return shp
    return None


def delete_table_row(table, row_idx):
    tbl = table._tbl
    trs = tbl.findall(qn('a:tr'))
    tbl.remove(trs[row_idx])


def set_cell_text(cell, text):
    """Replace the text of a cell while keeping the first run's formatting."""
    tf = cell.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        run = p.runs[0]
        run.text = str(text)
        # drop any extra runs left over in that paragraph
        for extra in p.runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        run = p.add_run()
        run.text = str(text)
    # drop any extra paragraphs beyond the first
    for extra_p in tf.paragraphs[1:]:
        extra_p._p.getparent().remove(extra_p._p)


def set_details_cell(cell, content_text, source_url):
    """
    Fills the 'Details' cell with the post content, then appends a
    hyperlinked word 'Source' at the end (only if a URL is available).
    """
    tf = cell.text_frame
    p = tf.paragraphs[0]
    run = p.runs[0] if p.runs else p.add_run()
    run.text = str(content_text)
    for extra in p.runs[1:]:
        extra._r.getparent().remove(extra._r)
    for extra_p in tf.paragraphs[1:]:
        extra_p._p.getparent().remove(extra_p._p)

    if source_url:
        sep = p.add_run()
        sep.text = "  "
        sep.font.size = run.font.size
        sep.font.name = run.font.name

        link_run = p.add_run()
        link_run.text = "Source"
        link_run.font.size = run.font.size
        link_run.font.name = run.font.name
        link_run.font.underline = True
        link_run.font.color.rgb = HYPERLINK_COLOR
        link_run.hyperlink.address = source_url


# ---------------------------------------------------------------------------
# 5. Main entry point
# ---------------------------------------------------------------------------
def generate_pptx(template_path, excel_file_like, rows_per_slide=None,
                   linkedin_id=None, exec_name=None, exec_country=None,
                   picture_bytes=None):
    """
    template_path : path to Sample.pptx
    excel_file_like : path or file-like object for the .xlsx upload
    rows_per_slide : how many table rows to put on each slide.
                      Defaults to however many data rows the template
                      slide already has (3 in Sample.pptx).
    linkedin_id, exec_name, exec_country : optional strings to stamp into
                      the header on every slide, so it's obvious whose
                      data the deck contains.
    picture_bytes   : optional raw image bytes (jpg/png) to use as the
                      executive's photo on every slide.
    Returns: BytesIO of the finished .pptx
    """
    records = read_excel_rows(excel_file_like)
    if not records:
        raise ValueError("No data rows found in the Excel file.")

    prs = Presentation(template_path)
    template_slide = prs.slides[0]
    template_table_shape = find_table_shape(template_slide)
    template_table = template_table_shape.table

    # Stamp the executive's branding onto the template slide BEFORE we
    # duplicate it, so every duplicated slide inherits the same photo/ID.
    set_branding(template_slide, linkedin_id, exec_name, exec_country, picture_bytes)

    n_body_rows = len(template_table.rows) - 1  # minus header row
    if rows_per_slide is None:
        rows_per_slide = n_body_rows

    chunks = [records[i:i + rows_per_slide] for i in range(0, len(records), rows_per_slide)]
    total_slides = len(chunks)

    # Fill the FIRST slide with chunk 0, using the existing template rows
    def fill_slide(slide, chunk, page_no):
        table_shape = find_table_shape(slide)
        table = table_shape.table

        # Make sure the table has exactly len(chunk) body rows
        current_body_rows = len(table.rows) - 1
        while current_body_rows > len(chunk):
            delete_table_row(table, current_body_rows)  # delete last body row
            current_body_rows -= 1
        # NOTE: if a chunk ever needs MORE rows than the template has,
        # increase rows_per_slide, or extend this function to clone a row.

        for i, rec in enumerate(chunk):
            row_idx = i + 1  # +1 to skip header row
            set_cell_text(table.cell(row_idx, 0), rec["post_type"])
            set_cell_text(table.cell(row_idx, 1), rec["timeline"])
            set_cell_text(table.cell(row_idx, 2), rec["headline"])
            set_details_cell(table.cell(row_idx, 3), rec["content"], rec["source_url"])

        # Update the "(1/2)" style page counter in the title, if present
        title_shape = find_title_shape(slide)
        if title_shape is not None:
            for p in title_shape.text_frame.paragraphs:
                for r in p.runs:
                    if re.search(r"\(\d+/\d+\)", r.text):
                        r.text = re.sub(r"\(\d+/\d+\)", f"({page_no}/{total_slides})", r.text)

    fill_slide(template_slide, chunks[0], 1)

    for page_no, chunk in enumerate(chunks[1:], start=2):
        new_slide = duplicate_slide(prs, 0)
        fill_slide(new_slide, chunk, page_no)

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out
